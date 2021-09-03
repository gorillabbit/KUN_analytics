import glob
import os
import textwrap
from pathlib import Path

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Pt

p = Path('H:/Projects/週刊KUN分析')  # フォルダを指定して色々取得しとく、動画の保存先も作成
folders = [str(x) for x in p.iterdir() if x.is_dir()][-1]
srcs = glob.glob(folders + '/*.xlsx')  # エクセルファイルからデータ取得

background_path = 'H:/Projects/週刊KUN分析/background.png'
mplus_bold = 'H:/Projects/週刊KUN分析/mplus-2m-bold.ttf'
mplus_black = 'H:/Projects/週刊KUN分析/mplus-2c-black.ttf'
fontcolor = '0x222222'

df_trans_this_week = pd.read_excel(srcs[0])
df_trans_last_week = pd.read_excel(srcs[1])
df_video_this_week = pd.read_excel(srcs[2])
df_video_last_week = pd.read_excel(srcs[3])
df_weekly = pd.read_excel(srcs[4])

slide_width = 2400
slide_height = 1350
weekly_ppt = Presentation()
weekly_ppt.slide_width = Pt(slide_width)
weekly_ppt.slide_height = Pt(slide_height)


def add_slide(ppt):
    slide_blank = ppt.slide_layouts[6]
    slide = ppt.slides.add_slide(slide_blank)
    slide.shapes.add_picture(background_path, 0, 0, width=Pt(slide_width), height=Pt(slide_height))
    return slide


def make_table(slide_name, row_q, x, y, col_1_w, col_2_w, col_1_title, col_2_title, col_q=2, col_3_w=0, style='blue'):
    table = slide_name.shapes.add_table(row_q + 1, col_q, Pt(x), Pt(y), Pt(0), Pt(0)).table
    table.columns[0].width = Pt(col_1_w)
    table.columns[1].width = Pt(col_2_w)
    table.cell(0, 0).text = col_1_title
    table.cell(0, 1).text = col_2_title
    if col_3_w != 0:
        table.columns[2].width = Pt(col_3_w)
    if style == 'red':
        tbl = table._graphic_frame.element.graphic.graphicData.tbl
        style_id = '{21E4AEA4-8DFA-4A89-87EB-49C32662AFE0}' #UUIDはgithub参照
        tbl[0][-1].text = style_id
    return table


def change_cell_font_and_size(cell, size, font='M+ 2p medium'):
    cell.text_frame.paragraphs[0].font.size = Pt(size)
    cell.text_frame.paragraphs[0].font.name = font


def change_text_font_and_size(text, font, size):
    text.text_frame.paragraphs[0].font.size = Pt(size)
    text.text_frame.paragraphs[0].font.name = font


def make_text(slide, x, y, inner_text, font_size=35, font='M+ 2p heavy'):
    text = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(100), Pt(50))
    text.text_frame.text = inner_text
    change_text_font_and_size(text, font, font_size)


def make_box(slide, x, y, w, h):
    box_1 = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    box_1.line.color.rgb = RGBColor(0, 0, 0)
    box_1.line.width = Pt(5)


def make_line(slide, start_x, start_y, end_x):
    line_1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(start_x), Pt(start_y), Pt(end_x), Pt(start_y))
    line_1.line.color.rgb = RGBColor(0, 0, 0)
    line_1.line.width = Pt(3)


font_text = 'M+ 2p medium'
font_title = 'M+ 2p heavy'
# 動画一覧
row_num = len(df_video_this_week)
slide = add_slide(weekly_ppt)
itiran = make_table(slide, row_num, 50, 50, 100, 1200, 'No', 'タイトル')
change_cell_font_and_size(itiran.cell(0, 0), 25)
change_cell_font_and_size(itiran.cell(0, 1), 25)
for i in range(row_num):
    row = df_video_this_week['タイトル'][i]
    itiran.cell(i+1, 0).text = str(i+1)
    change_cell_font_and_size(itiran.cell(i+1, 0), 18)
    itiran.cell(i+1, 1).text = row
    change_cell_font_and_size(itiran.cell(i+1, 1), 18)

# 指標一覧
index = make_table(slide, 20, 1400, 50, 400, 400, '指標', '')
change_cell_font_and_size(index.cell(0, 0), 35)
col_list = [0, 1, 2, 4, 5, 6, 7, 8, 9, 17, 10, 11, 18, 20, 12, 13, 19, 14, 15, 16]
for i, col in enumerate(col_list):
    row = str(df_weekly.iloc[-2, col])
    index.cell(i+1, 1).text = row
    change_cell_font_and_size(index.cell(i+1, 1), 30)
    index.cell(i+1, 0).text = df_weekly.columns[col]
    change_cell_font_and_size(index.cell(i+1, 0), 30)


# 週刊情報
def make_rank_table(rank_table, rank_df, text_long=100, text_size=30, title_size=35, row_count=5):
    change_cell_font_and_size(rank_table.cell(0, 0), title_size)
    change_cell_font_and_size(rank_table.cell(0, 1), title_size)
    for j in range(row_count):
        rank_table.cell(j+1, 0).text = str(j+1)
        change_cell_font_and_size(rank_table.cell(j+1, 0), text_size)
        title = textwrap.wrap(rank_df.iloc[j, 0], text_long)
        rank_table.cell(j+1, 1).text = title[0]
        change_cell_font_and_size(rank_table.cell(j+1, 1), text_size)
        rank_table.cell(j+1, 2).text = str(round(rank_df.iloc[j, 1],2))
        change_cell_font_and_size(rank_table.cell(j+1, 2), text_size)


weekly_graphs = glob.glob(folders + '/weekly/*.png')
video_num_slide = add_slide(weekly_ppt)  # 動画数
video_num_slide.shapes.add_picture(weekly_graphs[0], Pt(50), Pt(50), width=Pt(2350), height=Pt(1237))

weekly_slide = []
for i in range(1, 7):
    weekly_slide.append(add_slide(weekly_ppt))  # 再生数
    weekly_slide[i-1].shapes.add_picture(weekly_graphs[i], Pt(25), Pt(25), width=Pt(2350), height=Pt(900))

col_list = ['再生数', '高評価-低評価比率', '高評価数', '高評価数(再生数比x1000)', '低評価数', '低評価数(再生数比x1000)',
            'コメント数', 'コメント数(再生数比x1000)', '長さ(分)']
des_rank_df = []
asc_rank_df = []
for col in col_list:
    des_rank_df.append(df_video_this_week.sort_values(col, ascending=False)[['タイトル', col]])
    asc_rank_df.append(df_video_this_week.sort_values(col, ascending=True)[['タイトル', col]])

view_c_rank_table = make_table(weekly_slide[0], 5, 50, 1000, 90, 850, 'Rank', '再生数の順位(上から5本)', col_q=3, col_3_w=150)
view_c_rank2_table = make_table(weekly_slide[0], 5, 1250, 1000, 90, 850, 'Rank', '再生数の順位(下から5本)', col_q=3, col_3_w=150, style='red')
make_rank_table(view_c_rank_table, des_rank_df[0], text_long=30, text_size=20, title_size=23)
make_rank_table(view_c_rank2_table, asc_rank_df[0], text_long=30, text_size=20, title_size=23)
like_c_rank_table = make_table(weekly_slide[1], 5, 50, 930, 50, 600, 'R', '高評価数の順位(上から5本)', col_q=3, col_3_w=80)
unlike_c_rank_table = make_table(weekly_slide[1], 5, 800, 930, 50, 600, 'R', '低評価数の順位(上から5本)', col_q=3, col_3_w=80)
like_unlike_r_rank_table = make_table(weekly_slide[1], 5, 1550, 930, 50, 600, 'R', '高-低評価比率の順位(上から5本)', col_q=3, col_3_w=100)
like_c_rank2_table = make_table(weekly_slide[1], 5, 50, 1140, 50, 600, 'R', '高評価数の順位(下から5本)', col_q=3, col_3_w=80, style='red')
unlike_c_rank2_table = make_table(weekly_slide[1], 5, 800, 1140, 50, 600, 'R', '低評価数の順位(下から5本)', col_q=3, col_3_w=80, style='red')
like_unlike_r_rank2_table = make_table(weekly_slide[1], 5, 1550, 1140, 50, 600, 'R', '高-低評価比率の順位(下から5本)', col_q=3, col_3_w=100, style='red')
make_rank_table(like_c_rank_table, des_rank_df[2], text_long=27, text_size=20, title_size=23)
make_rank_table(unlike_c_rank_table, des_rank_df[4], text_long=27, text_size=20, title_size=23)
make_rank_table(like_unlike_r_rank_table, des_rank_df[1], text_long=27, text_size=20, title_size=23)
make_rank_table(like_c_rank2_table, asc_rank_df[2], text_long=27, text_size=20, title_size=23)
make_rank_table(unlike_c_rank2_table, asc_rank_df[4], text_long=27, text_size=20, title_size=23)
make_rank_table(like_unlike_r_rank2_table, asc_rank_df[1], text_long=27, text_size=20, title_size=23)


k = 1
for i, name in enumerate(['高評価数', '低評価数', 'コメント数']):
    count_rank_table = make_table(weekly_slide[i+2], 5, 50, 930, 80, 900, 'Rank', name+'のランキング(上から5本)', col_q=3, col_3_w=100)
    count_rank2_table = make_table(weekly_slide[i+2], 5, 50, 1140, 80, 900, 'Rank', name+'のランキング(下から5本)', col_q=3, col_3_w=100, style='red')
    count_per_rank_table = make_table(weekly_slide[i+2], 5, 1300, 930, 80, 900, 'Rank', '1000再生あたりの'+name+'のランキング(上から5本)', col_q=3, col_3_w=100)
    count_per_rank2_table = make_table(weekly_slide[i+2], 5, 1300, 1140, 80, 900, 'Rank', '1000再生あたりの'+name+'のランキング(下から5本)', col_q=3, col_3_w=100, style='red')
    k += 1
    make_rank_table(count_rank_table, des_rank_df[k], text_long=45, text_size=20, title_size=23)
    make_rank_table(count_rank2_table, asc_rank_df[k], text_long=45, text_size=20, title_size=23)
    k += 1
    make_rank_table(count_per_rank_table, des_rank_df[k], text_long=45, text_size=20, title_size=23)
    make_rank_table(count_per_rank2_table, asc_rank_df[k], text_long=45, text_size=20, title_size=23)


duration_rank_table = make_table(weekly_slide[5], 5, 50, 1000, 90, 850,'Rank', '動画の長さのランキング(上から5本)', col_q=3, col_3_w=120)
duration_rank2_table = make_table(weekly_slide[5], 5, 1250, 1000, 90, 850,'Rank', '動画の長さのランキング(下から5本)', col_q=3, col_3_w=100, style='red')
make_rank_table(duration_rank_table, des_rank_df[8], text_long=35, text_size=25, title_size=30)
make_rank_table(duration_rank2_table, asc_rank_df[8], text_long=35, text_size=25, title_size=30)

# 伸びのランキング(今週)


def make_nobi_rank_slide(nobi_list, slide, video_count, option, base_df):
    nobi_table_x = 50
    text_long = 45
    col_2_width = 800
    for j, title_nobi in enumerate(nobi_list):
        if j == 1:
            nobi_table_x = 1000
            text_long = 30
            col_2_width = 550
        elif j == 2:
            nobi_table_x = 1700
        nobi_table = make_table(slide, video_count, nobi_table_x, 50, 70, col_2_width, 'Rank', option+'投稿動画の'+title_nobi+'のランキング', col_q=3, col_3_w=60)
        nobi_table.cell(0, 2).text = '伸び'
        nobi_rank = base_df.sort_values(title_nobi, ascending=True)[['タイトル', title_nobi]]
        nobi_rank = nobi_rank[nobi_rank[title_nobi] != 0]
        make_rank_table(nobi_table, nobi_rank, text_long, 18, 20, row_count=len(nobi_rank))


nobi_slide = add_slide(weekly_ppt)
make_nobi_rank_slide(['伸び12時間', '伸び48時間', '伸び96時間'], nobi_slide, row_num, '今週', df_video_this_week)
nobi_last_slide = add_slide(weekly_ppt)
make_nobi_rank_slide(['伸び12時間', '伸び48時間', '伸び96時間'], nobi_last_slide, len(df_video_last_week), '先週', df_video_last_week)
nobi_last_slide_2 = add_slide(weekly_ppt)
make_nobi_rank_slide(['伸び144時間', '伸び192時間', '伸び240時間'], nobi_last_slide_2, len(df_video_last_week), '先週', df_video_last_week)
weekly_ppt.save(folders + '/weekly.pptx')


# 動画ごとのスライド
def make_daily_pptx(option, df, df_rank):
    daily_ppt = Presentation()
    daily_ppt.slide_width = Pt(slide_width)
    daily_ppt.slide_height = Pt(slide_height)

    daily_graphs = glob.glob(folders + '/vct_day' + option + '/*.png')
    video_graphs = glob.glob(folders + '/vct_video' + option + '/*.png')
    thumbnails = glob.glob(folders + '/thumbnail' + option + '/*.png')

    f_s_s = 35
    f_s_l = 65

    for daily_graph in daily_graphs:
        daily_slide = add_slide(daily_ppt)
        daily_slide.shapes.add_picture(daily_graph, Pt(100), Pt(50), width=Pt(2200), height=Pt(1237))
        make_box(daily_slide, 100, 50, 2200, 1237)
        date = os.path.splitext(os.path.basename(daily_graph))[0]
        index_list = [['再生数'], ['高評価数', '高評価数(再生数比x1000)'], ['低評価数', '低評価数(再生数比x1000)'], ['コメント数', 'コメント数(再生数比x1000)']]
        nobi_list = ['伸び12時間', '伸び48時間','伸び96時間', '伸び144時間', '伸び192時間', '伸び240時間']
        date_list = ['(半日)', '(2日)', '(4日)', '(6日)', '(8日)', '(10日)']
        print(date)
        df_rank_pct = df.replace(0, np.nan)
        df_rank_pct = df_rank_pct.rank(numeric_only=True, ascending=True, method='min', pct=True)
        for video_id in df[df['投稿日'] == date]['videoID']:
            video = df[df['videoID'] == video_id].iloc[0]
            video_rank = df_rank[df['videoID'] == video_id].iloc[0]
            video_rank_pct = df_rank_pct[df['videoID'] == video_id].iloc[0]
            print(video_id)
            video_slide = add_slide(daily_ppt)
            thumbnail = [s for s in thumbnails if video_id in s][0]

            def shape_slide(image):
                video_slide.shapes.add_picture(image, Pt(50), Pt(50), width=Pt(1600), height=Pt(900))
                make_box(video_slide, 50, 50, 1600, 900)
                make_box(video_slide, 50, 1000, 2300, 300)
                make_box(video_slide, 1700, 50, 650, 900)
                title = textwrap.wrap(video['タイトル'], 45)  # 長いタイトルを分割
                title.append('')  # タイトルが短い場合にtitle[1]が無くなるの回避
                make_text(video_slide, 70, 1030, title[0], font_size=50)
                make_text(video_slide, 1200, 1080, title[1], font_size=50)
                make_text(video_slide, 100, 1100, '投稿日時 ' + video['投稿日時'])
                make_text(video_slide, 100, 1150, '動画ID　' + video_id)
                make_text(video_slide, 100, 1200, 'シリーズ ' + str(video.iloc[11]))

                gap_count = 0
                for count, index in enumerate(index_list):
                    place_title = f_s_s*(count)+f_s_l*count+5*gap_count
                    gap_count += 1
                    place = f_s_s*(count+1)+f_s_l*count+5*gap_count
                    gap_count += 2

                    make_text(video_slide, 1720, 80+place_title, index[0])
                    make_text(video_slide, 1810, 80+place, str(video[index[0]]), font_size=f_s_l)
                    make_text(video_slide, 1750, 97+place, str(video_rank[index[0]])[:2])
                    make_line(video_slide, 1727, 85+place, 1850)

                    if len(index) == 2:
                        make_text(video_slide, 1750 + 300, 87+place_title, '再生数比x1000', font_size=30)
                        make_text(video_slide, 1750 + 370, 85+place, str(round(video[index[1]], 2)), font_size=55)
                        make_text(video_slide, 1750 + 300, 97+place, str(video_rank[index[1]])[:2])
                        make_line(video_slide, 1850, 85+place, 2280)

                make_text(video_slide, 1720, 80+f_s_s*4+f_s_l*4+5*12, '動画の長さ')
                make_text(video_slide, 1810, 80+f_s_s*5+f_s_l*4+5*13, str(video['長さ']), font_size=f_s_l)
                make_text(video_slide, 1750, 97+f_s_s*5+f_s_l*4+5*13, str(video_rank['長さ(分)'])[:2])
                make_line(video_slide, 1727, 80+f_s_s*5+f_s_l*4+5*14, 1900)

                make_text(video_slide, 1720, 80+f_s_s*5+f_s_l*5+5*15, '動画の勢い(直近300本中順位)')
                make_line(video_slide, 1727, 80+f_s_s*6+f_s_l*5+5*17, 2220)
                if option == '':
                    for count in range(3):
                        offset = f_s_s*6+f_s_l*(5+count)+5*16
                        if video[nobi_list[count]] != 0:  # 伸びが0以外のとき記入
                            make_text(video_slide, 1750, 97+offset, str(int(video_rank_pct[nobi_list[count]]*100))+'%')
                        make_text(video_slide, 1750+100, 80+offset, str(video[nobi_list[count]]), font_size=f_s_l)
                        make_text(video_slide, 1750+100+140, 90+offset, date_list[count])

                elif option == '_last':
                    for count in range(3):
                        offset = f_s_s*6+f_s_l*(5+count)+5*16
                        if video[nobi_list[count]] != 0:
                            make_text(video_slide, 1750, 90+offset, str(int(video_rank_pct[nobi_list[count]]*100))+'%', font_size=30)
                        make_text(video_slide, 1750+85, 80+offset, str(video[nobi_list[count]]), font_size=50)
                        make_text(video_slide, 1750+85+100, 90+offset, date_list[count])

                        if video[nobi_list[count+3]] != 0:
                            make_text(video_slide, 2040, 90+offset, str(int(video_rank_pct[nobi_list[count+3]]*100))+'%', font_size=30)
                        make_text(video_slide, 2040+85, 80+offset, str(video[nobi_list[count+3]]), font_size=50)
                        make_text(video_slide, 2040+85+100, 90+offset, date_list[count+3])

            shape_slide(thumbnail)
            video_slide = add_slide(daily_ppt)
            graph = [s for s in video_graphs if video_id in s][0]
            shape_slide(graph)

    add_slide(daily_ppt)
    daily_ppt.save(folders + '/daily' + option + '.pptx')


df_rank_this = df_video_this_week.rank(numeric_only=True, ascending=False, method='dense')
make_daily_pptx('', df_video_this_week, df_rank_this)
df_rank_last = df_video_last_week.rank(numeric_only=True, ascending=False, method='dense')
make_daily_pptx('_last', df_video_last_week, df_rank_last)

for i in range(3):
    opening_ppt = Presentation()
    opening_ppt.slide_width = Pt(slide_width)
    opening_ppt.slide_height = Pt(slide_height)
    thumbnail = 'H:/Projects/週刊KUN分析/thumbnail_'+str(i)+'.png'
    opening_slide = add_slide(opening_ppt)

    opening_slide.shapes.add_picture(thumbnail, Pt(0), Pt(0), width=Pt(slide_width), height=Pt(slide_height))
    make_text(opening_slide, 1000, 900, str(df_weekly.iloc[-2, 21]), font_size=200, font='M+ 2p black')
    opening_ppt.save(folders + '/opening_slide_'+str(i)+'.pptx')
